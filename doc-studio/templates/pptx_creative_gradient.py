#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
クリエイティブグラデーションテーマ PPTX
- カラー: パープルグラデーション(#667eea→#764ba2) + ピンクアクセント
- 用途: クリエイティブエージェンシー紹介
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
PURPLE_1 = RGBColor(0x66, 0x7e, 0xea)
PURPLE_2 = RGBColor(0x76, 0x4b, 0xa2)
PINK = RGBColor(0xf0, 0x93, 0xfb)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x2d, 0x2d, 0x2d)


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_gradient_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # パープル背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE_2
    bg.line.fill.background()

    # 円形装飾1
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PURPLE_1
    circle1.fill.fore_color.brightness = 0.2
    circle1.line.fill.background()

    # 円形装飾2
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = PINK
    circle2.fill.fore_color.brightness = 0.3
    circle2.line.fill.background()

    # メインタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_creative_content_slide(prs, title, content_items, layout_type="bullets"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # 左側カラーバー
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = PURPLE_1
    left_bar.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_2

    # ピンクアクセント
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.4), Inches(0.15), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PINK
    accent.line.fill.background()

    # コンテンツ
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"◆ {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(20)

    return slide

def add_full_image_slide(prs, title, description):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # パープル背景（画像風）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE_1
    bg.line.fill.background()

    # グラデーションオーバーレイ風
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4), Inches(13.333), Inches(3.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PURPLE_2
    overlay.fill.fore_color.brightness = -0.2
    overlay.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 説明
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = PINK

    return slide

def create_presentation(output_path: str, _data: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトル
    add_gradient_title_slide(prs, "CREATIVE STUDIO", "Design that inspires")

    # 2. コンテンツ
    add_creative_content_slide(prs, "Our Philosophy", [
        "デザインは問題解決の手段である",
        "ユーザー中心の思考を徹底する",
        "技術とアートの融合を追求する",
        "持続可能な創造性を大切にする"
    ])

    # 3. フルイメージ
    add_full_image_slide(prs, "Brand Identity",
        "あなたのブランドの本質を視覚化し、\n世界に届けるデザインを創造します")

    # 4. サービス
    add_creative_content_slide(prs, "Services", [
        "Branding & Identity: ブランド戦略からビジュアルまで",
        "Web Design: レスポンシブで使いやすいWebサイト",
        "UI/UX Design: 直感的なユーザー体験の設計",
        "Motion Graphics: 動きのある視覚表現"
    ])

    # 5. プロセス
    add_creative_content_slide(prs, "Our Process", [
        "Discover: 課題の発見とリサーチ",
        "Define: 戦略の定義と方向性の設定",
        "Design: クリエイティブな発想とデザイン",
        "Deliver: 完成度の高い成果物の納品"
    ])

    # 6. フルイメージ
    add_full_image_slide(prs, "Let's Create Together",
        "次のプロジェクトで一緒に素晴らしいものを\n作り上げませんか？")

    # 7. コンタクト
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE_2
    bg.line.fill.background()

    # 円装飾
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2), Inches(3.333), Inches(3.333))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PINK
    circle.fill.fore_color.brightness = 0.2
    circle.line.fill.background()

    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Get in Touch"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate creative gradient PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/02_creative_gradient.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_presentation(args.output, data)
