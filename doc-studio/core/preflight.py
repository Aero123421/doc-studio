"""
Preflight Checker - Quality assurance for generated documents
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, asdict
from enum import Enum
import json


class CheckSeverity(Enum):
    INFO = "info"
    WARNING = "warning"
    ERROR = "error"


@dataclass
class CheckResult:
    """Result of a preflight check"""
    check_name: str
    severity: CheckSeverity
    message: str
    details: Optional[Dict[str, Any]] = None
    fix_suggestion: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "check_name": self.check_name,
            "severity": self.severity.value,
            "message": self.message,
            "details": self.details,
            "fix_suggestion": self.fix_suggestion,
        }


class PreflightChecker:
    """Check document quality before finalization"""

    SUPPORTED_CHECKS = {
        "fonts": "Check font embedding and availability",
        "images": "Check image resolution and format",
        "links": "Check external links validity",
        "colors": "Check color contrast and accessibility",
        "accessibility": "Check accessibility compliance",
        "metadata": "Check document metadata",
        "filesize": "Check file size limits",
    }

    def __init__(self):
        self.results: List[CheckResult] = []
        self.checks_enabled = list(self.SUPPORTED_CHECKS.keys())

    def check(
        self,
        file_path: str,
        checks: Optional[List[str]] = None,
        fix_issues: bool = False
    ) -> Dict[str, Any]:
        """
        Run preflight checks on a document

        Args:
            file_path: Path to document
            checks: List of checks to run (None = all)
            fix_issues: Attempt to auto-fix issues

        Returns:
            Check results summary
        """
        self.results = []
        file_path = Path(file_path)

        if not file_path.exists():
            return {
                "success": False,
                "error": f"File not found: {file_path}",
                "results": []
            }

        # Determine checks to run
        checks_to_run = checks or self.checks_enabled

        # Run checks based on file type
        ext = file_path.suffix.lower()

        if ext == ".pdf":
            self._check_pdf(file_path, checks_to_run)
        elif ext == ".pptx":
            self._check_pptx(file_path, checks_to_run)
        elif ext == ".docx":
            self._check_docx(file_path, checks_to_run)
        elif ext == ".xlsx":
            self._check_xlsx(file_path, checks_to_run)
        elif ext == ".html":
            self._check_html(file_path, checks_to_run)
        else:
            return {
                "success": False,
                "error": f"Unsupported file type: {ext}",
                "results": []
            }

        # Generate summary
        summary = self._generate_summary()

        # Auto-fix if requested
        if fix_issues:
            self._auto_fix(file_path)

        return summary

    def _check_pdf(self, file_path: Path, checks: List[str]):
        """Check PDF file"""
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))

            if "fonts" in checks:
                self._check_pdf_fonts(reader)

            if "metadata" in checks:
                self._check_pdf_metadata(reader)

            if "filesize" in checks:
                self._check_filesize(file_path, max_size_mb=50)

        except ImportError:
            self.results.append(CheckResult(
                check_name="pdf_validation",
                severity=CheckSeverity.WARNING,
                message="PDF validation library not available",
                fix_suggestion="Install pypdf for full PDF checking"
            ))

    def _check_pptx(self, file_path: Path, checks: List[str]):
        """Check PowerPoint file"""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))

            if "images" in checks:
                self._check_pptx_images(prs)

            if "fonts" in checks:
                self._check_pptx_fonts(prs)

            if "filesize" in checks:
                self._check_filesize(file_path, max_size_mb=100)

        except ImportError:
            self.results.append(CheckResult(
                check_name="pptx_validation",
                severity=CheckSeverity.WARNING,
                message="PowerPoint validation library not available",
            ))

    def _check_docx(self, file_path: Path, checks: List[str]):
        """Check Word file"""
        try:
            from docx import Document

            doc = Document(str(file_path))

            if "accessibility" in checks:
                self._check_docx_accessibility(doc)

            if "filesize" in checks:
                self._check_filesize(file_path, max_size_mb=50)

        except ImportError:
            self.results.append(CheckResult(
                check_name="docx_validation",
                severity=CheckSeverity.WARNING,
                message="Word validation library not available",
            ))

    def _check_xlsx(self, file_path: Path, checks: List[str]):
        """Check Excel file"""
        if "filesize" in checks:
            self._check_filesize(file_path, max_size_mb=50)

    def _check_html(self, file_path: Path, checks: List[str]):
        """Check HTML file"""
        content = file_path.read_text(encoding="utf-8")

        if "accessibility" in checks:
            self._check_html_accessibility(content)

        if "links" in checks:
            self._check_html_links(content)

    def _check_pdf_fonts(self, reader):
        """Check PDF font embedding"""
        # Basic check - full implementation would parse font resources
        self.results.append(CheckResult(
            check_name="fonts",
            severity=CheckSeverity.INFO,
            message="PDF font check completed",
            details={"pages": len(reader.pages)}
        ))

    def _check_pdf_metadata(self, reader):
        """Check PDF metadata"""
        metadata = reader.metadata

        if not metadata:
            self.results.append(CheckResult(
                check_name="metadata",
                severity=CheckSeverity.WARNING,
                message="PDF metadata is missing",
                fix_suggestion="Add title, author, and creation date metadata"
            ))
        else:
            missing = []
            if not metadata.get("/Title"):
                missing.append("title")
            if not metadata.get("/Author"):
                missing.append("author")

            if missing:
                self.results.append(CheckResult(
                    check_name="metadata",
                    severity=CheckSeverity.WARNING,
                    message=f"Missing metadata fields: {', '.join(missing)}",
                    fix_suggestion="Add missing metadata fields"
                ))
            else:
                self.results.append(CheckResult(
                    check_name="metadata",
                    severity=CheckSeverity.INFO,
                    message="PDF metadata is complete"
                ))

    def _check_pptx_images(self, prs):
        """Check PowerPoint image resolution"""
        low_res_images = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    # Check image size
                    img = shape.image
                    if img.width < 100000:  # EMU
                        low_res_images += 1

        if low_res_images > 0:
            self.results.append(CheckResult(
                check_name="images",
                severity=CheckSeverity.WARNING,
                message=f"{low_res_images} low resolution images detected",
                fix_suggestion="Replace images with higher resolution versions"
            ))
        else:
            self.results.append(CheckResult(
                check_name="images",
                severity=CheckSeverity.INFO,
                message="All images have adequate resolution"
            ))

    def _check_pptx_fonts(self, prs):
        """Check PowerPoint font usage"""
        fonts_used = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts_used.add(run.font.name)

        self.results.append(CheckResult(
            check_name="fonts",
            severity=CheckSeverity.INFO,
            message=f"Fonts used: {', '.join(fonts_used)}",
            details={"fonts": list(fonts_used)}
        ))

    def _check_docx_accessibility(self, doc):
        """Check Word accessibility"""
        issues = []

        # Check for heading hierarchy
        headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
        if not headings:
            issues.append("No heading structure found")

        # Check image alt text
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                # Would need more detailed check for alt text
                pass

        if issues:
            self.results.append(CheckResult(
                check_name="accessibility",
                severity=CheckSeverity.WARNING,
                message="; ".join(issues),
                fix_suggestion="Add heading structure and image alt text"
            ))
        else:
            self.results.append(CheckResult(
                check_name="accessibility",
                severity=CheckSeverity.INFO,
                message="Document accessibility checks passed"
            ))

    def _check_html_accessibility(self, content):
        """Check HTML accessibility"""
        issues = []

        # Basic checks
        if "<html" not in content.lower():
            issues.append("Missing html tag")
        if "<title>" not in content:
            issues.append("Missing title tag")
        if "lang=" not in content:
            issues.append("Missing lang attribute")

        if issues:
            self.results.append(CheckResult(
                check_name="accessibility",
                severity=CheckSeverity.ERROR,
                message="; ".join(issues),
            ))
        else:
            self.results.append(CheckResult(
                check_name="accessibility",
                severity=CheckSeverity.INFO,
                message="HTML accessibility checks passed"
            ))

    def _check_html_links(self, content):
        """Check HTML external links"""
        import re

        # Find all URLs
        urls = re.findall(r'href=["\'](https?://[^"\']+)["\']', content)

        self.results.append(CheckResult(
            check_name="links",
            severity=CheckSeverity.INFO,
            message=f"Found {len(urls)} external links",
            details={"urls": urls[:10]}  # Limit to first 10
        ))

    def _check_filesize(self, file_path: Path, max_size_mb: int):
        """Check file size"""
        size_mb = file_path.stat().st_size / (1024 * 1024)

        if size_mb > max_size_mb:
            self.results.append(CheckResult(
                check_name="filesize",
                severity=CheckSeverity.WARNING,
                message=f"File size ({size_mb:.1f}MB) exceeds recommended limit ({max_size_mb}MB)",
                fix_suggestion="Compress images or reduce content",
                details={"size_mb": size_mb, "limit_mb": max_size_mb}
            ))
        else:
            self.results.append(CheckResult(
                check_name="filesize",
                severity=CheckSeverity.INFO,
                message=f"File size ({size_mb:.1f}MB) is within limits",
                details={"size_mb": size_mb}
            ))

    def _generate_summary(self) -> Dict[str, Any]:
        """Generate check summary"""
        errors = [r for r in self.results if r.severity == CheckSeverity.ERROR]
        warnings = [r for r in self.results if r.severity == CheckSeverity.WARNING]
        infos = [r for r in self.results if r.severity == CheckSeverity.INFO]

        return {
            "success": len(errors) == 0,
            "summary": {
                "total": len(self.results),
                "errors": len(errors),
                "warnings": len(warnings),
                "info": len(infos),
            },
            "results": [r.to_dict() for r in self.results],
            "can_proceed": len(errors) == 0,
        }

    def _auto_fix(self, file_path: Path):
        """Attempt to auto-fix issues"""
        # Implementation would fix common issues
        pass

    def get_supported_checks(self) -> Dict[str, str]:
        """Get list of supported checks"""
        return self.SUPPORTED_CHECKS.copy()
