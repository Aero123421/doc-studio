#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ビジネスモダンテーマ PPTX
- カラー: ネイビー(#1e3a5f) + ゴールド(#c9a227) + ホワイト
- 用途: デジタルトランスフォーメーション戦略
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
NAVY = RGBColor(0x1e, 0x3a, 0x5f)
GOLD = RGBColor(0xc9, 0xa2, 0x27)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)

    # 背景（ネイビーグラデーション風）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # ゴールドライン（上）
    line_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = GOLD
    line_top.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # ゴールドライン（下）
    line_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.5), Inches(5.333), Inches(0.03))
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = GOLD
    line_bottom.line.fill.background()

    return slide

def add_section_slide(prs, number, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # 左側ネイビーバー
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = NAVY
    left_bar.line.fill.background()

    # セクション番号
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(12), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{number}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # ヘッダーバー
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # タイトル（ヘッダー内）
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ゴールドアクセントライン
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(1), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # コンテンツ
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▶ {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(16)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # ヘッダー
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 左カラム
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5))
    tf = left_content_box.text_frame
    tf.word_wrap = True
    for item in left_content:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # センターライン
    center_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(0.02), Inches(5))
    center_line.fill.solid()
    center_line.fill.fore_color.rgb = LIGHT_GRAY
    center_line.line.fill.background()

    # 右カラム
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.8), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    right_content_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.8), Inches(4.5))
    tf = right_content_box.text_frame
    tf.word_wrap = True
    for item in right_content:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    return slide

def add_chart_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()

    # ヘッダー
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # バーチャート風の図形
    values = [0.7, 0.85, 0.6, 0.9, 0.75]
    labels = ["2020", "2021", "2022", "2023", "2024"]
    colors = [NAVY, GOLD, NAVY, GOLD, NAVY]

    for i, (val, label, color) in enumerate(zip(values, labels, colors)):
        x = Inches(1.5 + i * 2.2)
        height = Inches(val * 4)
        y = Inches(5.5 - val * 4)

        # バー
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.5), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # ラベル
        label_box = slide.shapes.add_textbox(x, Inches(5.6), Inches(1.5), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        # 値
        value_box = slide.shapes.add_textbox(x, y - Inches(0.4), Inches(1.5), Inches(0.4))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{int(val*100)}%"
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation(output_path: str, _data: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトル
    add_title_slide(prs,
        "Digital Transformation Strategy",
        "2025年度 事業戦略プレゼンテーション")

    # 2. セクション1
    add_section_slide(prs, 1, "Executive Summary")

    # 3. コンテンツ
    add_content_slide(prs, "経営方針", [
        "デジタル技術を活用した業務効率化と新規事業創出",
        "顧客体験（CX）の向上を第一優先に据えた投資計画",
        "データドリブン経営の徹底的な推進",
        "グリーンIT initiatives による持続可能な社会貢献"
    ])

    # 4. 2カラム
    add_two_column_slide(prs, "現状分析",
        "強み (Strengths)", [
            "豊富な顧客基盤（累計100社以上）",
            "高い技術力と専門人材",
            "安定的な財務基盤"
        ],
        "課題 (Challenges)", [
            "レガシーシステムの老朽化",
            "デジタル人材の不足",
            "競合他社との差別化"
        ])

    # 5. セクション2
    add_section_slide(prs, 2, "Growth Strategy")

    # 6. チャート
    add_chart_slide(prs, "売上成長推移")

    # 7. コンテンツ
    add_content_slide(prs, "戦略的投資領域", [
        "AI/機械学習基盤の構築（投資額: 5億円）",
        "クラウドネイティブ基盤への移行",
        "サイバーセキュリティ強化",
        "人材育成プログラムの拡充"
    ])

    # 8. 最終
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ご質問・ご相談はこちらまで"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate business modern PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/01_business_modern.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_presentation(args.output, data)
