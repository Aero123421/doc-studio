#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
テクニカルダークテーマ PPTX
- カラー: ダーク(#0d1117) + シアン(#00d4aa) + オレンジ(#ff6b35)
- 用途: システムアーキテクチャ解説
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# カラーパレット
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
DARK_CARD = RGBColor(0x16, 0x1b, 0x22)
CYAN = RGBColor(0x00, 0xd4, 0xaa)
ORANGE = RGBColor(0xff, 0x6b, 0x35)
WHITE = RGBColor(0xe6, 0xed, 0xf3)
GRAY = RGBColor(0x8b, 0x94, 0x9e)


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_dark_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # ダーク背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # グリッド線風装飾
    for i in range(5):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 3), 0, Inches(0.01), Inches(7.5))
        line.fill.solid()
        line.fill.fore_color.rgb = DARK_CARD
        line.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = CYAN
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER

    # プロンプト風カーソル
    cursor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(5.5), Inches(0.1), Inches(0.4))
    cursor.fill.solid()
    cursor.fill.fore_color.rgb = CYAN
    cursor.line.fill.background()

    return slide

def add_code_slide(prs, title, code_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # ダーク背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"$ {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = "Consolas"

    # コードブロック背景
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = DARK_CARD
    code_bg.line.color.rgb = CYAN
    code_bg.line.width = Pt(1)

    # コード
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        if line.startswith('#'):
            p.font.color.rgb = GRAY
        elif 'def' in line or 'class' in line:
            p.font.color.rgb = CYAN
        elif '"' in line or "'" in line:
            p.font.color.rgb = ORANGE
        else:
            p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    return slide

def add_architecture_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # ダーク背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # アーキテクチャ図（シンプルなボックス構成）
    components = [
        ("Frontend\n(React)", 1, 1, CYAN),
        ("API Gateway\n(Kong)", 5, 1, ORANGE),
        ("Backend\n(Go)", 9, 1, CYAN),
        ("Database\n(PostgreSQL)", 5, 4, ORANGE),
    ]

    for text, x, y, color in components:
        # ボックス
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_CARD
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # テキスト
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.4), Inches(3), Inches(1))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 接続線
    lines = [
        (Inches(4), Inches(1.75), Inches(5), Inches(1.75)),
        (Inches(8), Inches(1.75), Inches(9), Inches(1.75)),
        (Inches(6.5), Inches(2.5), Inches(6.5), Inches(4)),
    ]
    for x1, y1, x2, y2 in lines:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = GRAY
        line.line.width = Pt(2)

    return slide

def create_presentation(output_path: str, _data: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトル
    add_dark_title_slide(prs, "System Architecture", "Microservices Design Patterns")

    # 2. コードスライド
    add_code_slide(prs, "microservices.yaml", [
        "# Kubernetes Deployment",
        "apiVersion: apps/v1",
        "kind: Deployment",
        "metadata:",
        "  name: api-gateway",
        "spec:",
        "  replicas: 3",
        "  selector:",
        "    matchLabels:",
        "      app: gateway",
        "  template:",
        "    spec:",
        "      containers:",
        "      - name: kong",
        "        image: kong:3.0"
    ])

    # 3. アーキテクチャ
    add_architecture_slide(prs, "System Overview")

    # 4. コードスライド2
    add_code_slide(prs, "service.go", [
        "package main",
        "",
        "import (",
        '    "context"',
        '    "log"',
        "    ",
        '    "github.com/gin-gonic/gin"',
        ")",
        "",
        "func main() {",
        "    r := gin.Default()",
        '    r.GET("/health", healthCheck)',
        "    r.Run(':8080')",
        "}"
    ])

    # 5. パフォーマンス指標
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Metrics"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # メトリクスカード
    metrics = [
        ("Latency", "12ms", CYAN),
        ("Throughput", "10K rps", ORANGE),
        ("Uptime", "99.99%", CYAN),
        ("Error Rate", "0.01%", ORANGE),
    ]

    for i, (label, value, color) in enumerate(metrics):
        x = Inches(0.5 + i * 3.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(2.8), Inches(3))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_CARD
        card.line.color.rgb = color
        card.line.width = Pt(1)

        label_box = slide.shapes.add_textbox(x, Inches(2.5), Inches(2.8), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        value_box = slide.shapes.add_textbox(x, Inches(3.3), Inches(2.8), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate technical dark PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/03_technical_dark.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_presentation(args.output, data)
