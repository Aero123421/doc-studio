#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ミニマリストテーマ PPTX
- カラー: 白 + 黒 + アクセント1色のみ
- 用途: ミニマリストライフスタイル
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY_LIGHT = RGBColor(0xf5, 0xf5, 0xf5)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xe8, 0x5d, 0x4d)  # テラコッタ


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_minimal_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # タイトル（中央寄せ、大きな余白）
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_minimal_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # タイトル（左上、小さめ）
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK

    # 細いライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(0.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # コンテンツ（シンプル）
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY_TEXT
        p.space_after = Pt(24)

    return slide

def add_quote_slide(prs, quote, author):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # ライトグレー背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRAY_LIGHT
    bg.line.fill.background()

    # クォートマーク
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1.5))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(72)
    p.font.color.rgb = ACCENT

    # 引用文
    quote_box = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(9), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK

    # 著者
    author_box = slide.shapes.add_textbox(Inches(3), Inches(4.8), Inches(9), Inches(0.6))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {author}"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_TEXT

    return slide

def create_presentation(output_path: str, _data: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトル
    add_minimal_title_slide(prs, "Less is More", "ミニマリストの生き方")

    # 2. コンテンツ
    add_minimal_content_slide(prs, "ミニマリズムとは", [
        "持たない幸せを見つける",
        "本当に大切なものだけを残す",
        "シンプルだからこそ豊かになる"
    ])

    # 3. 引用
    add_quote_slide(prs,
        "完璧であるとは、追加するものがないことではなく、\n削ぎ落とすものがないことである。",
        "Antoine de Saint-Exupéry")

    # 4. コンテンツ2
    add_minimal_content_slide(prs, "実践のステップ", [
        "所有物を半分にする",
        "デジタルデトックス",
        "心に余白を作る"
    ])

    # 5. 最終
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    end_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = end_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate minimalist PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/04_minimalist.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_presentation(args.output, data)
