#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
コーポレートフォーマルテーマ PPTX
- カラー: 濃紺(#002147) + グレー(#f5f5f5) + レッド(#c41e3a)
- 用途: 決算説明会
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
NAVY = RGBColor(0x00, 0x21, 0x47)
RED = RGBColor(0xc4, 0x1e, 0x3a)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY_LIGHT = RGBColor(0xf5, 0xf5, 0xf5)
GRAY_TEXT = RGBColor(0x5a, 0x5a, 0x5a)


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_formal_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # 上部濃紺バー
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()

    # 赤ライン
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2), Inches(13.333), Inches(0.08))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = RED
    red_line.line.fill.background()

    # ロゴスペース（左上）
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.8))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "COMPANY LOGO"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_TEXT

    # フッター情報
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "株式会社サンプル | 2025年3月期 決算説明会"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT

    return slide

def add_formal_content_slide(prs, title, content, page_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # ヘッダー
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
    tf = header_title.text_frame
    p = tf.paragraphs[0]
    p.text = "2025年3月期 決算説明会"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ページ番号
    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.3), Inches(0.8), Inches(0.6))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 赤アクセントライン
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(0.6), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    # コンテンツ
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"・{item}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.space_after = Pt(12)

    # フッター
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(13.333), Inches(0.2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()

    return slide

def add_financial_slide(prs, page_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # ヘッダー
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.6))
    tf = header_title.text_frame
    p = tf.paragraphs[0]
    p.text = "2025年3月期 決算説明会"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    page_box = slide.shapes.add_textbox(Inches(12), Inches(0.3), Inches(0.8), Inches(0.6))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "財務ハイライト"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 財務データ（シンプル表）
    data = [
        ("項目", "2024年3月期", "2025年3月期", "増減率"),
        ("売上高", "100億円", "120億円", "+20%"),
        ("営業利益", "10億円", "15億円", "+50%"),
        ("経常利益", "9億円", "14億円", "+55%"),
        ("純利益", "6億円", "10億円", "+66%"),
    ]

    for i, (col1, col2, col3, col4) in enumerate(data):
        y = Inches(2.3 + i * 0.8)
        bg_color = NAVY if i == 0 else (GRAY_LIGHT if i % 2 == 0 else WHITE)
        text_color = WHITE if i == 0 else GRAY_TEXT

        # 行背景
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12.333), Inches(0.7))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.fill.background()

        # 項目名
        col1_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.15), Inches(3), Inches(0.5))
        tf = col1_box.text_frame
        p = tf.paragraphs[0]
        p.text = col1
        p.font.size = Pt(12)
        p.font.bold = (i == 0)
        p.font.color.rgb = text_color

        # 前期
        col2_box = slide.shapes.add_textbox(Inches(4), y + Inches(0.15), Inches(2.5), Inches(0.5))
        tf = col2_box.text_frame
        p = tf.paragraphs[0]
        p.text = col2
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.RIGHT

        # 当期
        col3_box = slide.shapes.add_textbox(Inches(7), y + Inches(0.15), Inches(2.5), Inches(0.5))
        tf = col3_box.text_frame
        p = tf.paragraphs[0]
        p.text = col3
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY if i > 0 else text_color
        p.alignment = PP_ALIGN.RIGHT

        # 増減率
        col4_box = slide.shapes.add_textbox(Inches(10), y + Inches(0.15), Inches(2.5), Inches(0.5))
        tf = col4_box.text_frame
        p = tf.paragraphs[0]
        p.text = col4
        p.font.size = Pt(12)
        p.font.color.rgb = RED if i > 0 else text_color
        p.alignment = PP_ALIGN.RIGHT

    # フッター
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(13.333), Inches(0.2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()

    return slide

def create_presentation(output_path: str, _data: dict | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトル
    add_formal_title_slide(prs,
        "2025年3月期 決算説明会",
        "Fiscal Year 2025 Financial Results")

    # 2. 議題
    add_formal_content_slide(prs, "本日の議題", [
        "2025年3月期 業績概要",
        "セグメント別業績",
        "財務状況",
        "株主還元方針",
        "今後の事業戦略"
    ], 2)

    # 3. 業績概要
    add_formal_content_slide(prs, "業績概要", [
        "売上高：120億円（前年比+20%）",
        "営業利益：15億円（前年比+50%）",
        "経常利益：14億円（前年比+55%）",
        "四半期連続の増収増益を達成"
    ], 3)

    # 4. 財務ハイライト（表）
    add_financial_slide(prs, 4)

    # 5. フッター付き最終スライド
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "ご清聴ありがとうございました"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    contact = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = contact.text_frame
    p = tf.paragraphs[0]
    p.text = "投資者関係部門：ir@example.com"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_LIGHT
    p.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(13.333), Inches(0.2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RED
    footer.line.fill.background()

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate corporate formal PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/05_corporate_formal.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_presentation(args.output, data)
