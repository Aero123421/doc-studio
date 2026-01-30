#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
高度なビジネスPPTX - エンタープライズ品質
- 複雑な図形とアニメーション準備
- データ可視化
- インフォグラフィック
- プレミアムデザイン
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import math

# プレミアムカラーパレット
COLORS = {
    'primary': RGBColor(0x1a, 0x36, 0x5d),
    'secondary': RGBColor(0x2c, 0x52, 0x82),
    'accent': RGBColor(0xc9, 0xa2, 0x27),
    'accent2': RGBColor(0xe8, 0x5d, 0x4d),
    'dark': RGBColor(0x1a, 0x20, 0x2c),
    'gray': RGBColor(0x71, 0x80, 0x96),
    'light': RGBColor(0xf7, 0xfa, 0xfc),
    'white': RGBColor(0xff, 0xff, 0xff),
    'success': RGBColor(0x48, 0xbb, 0x78),
    'warning': RGBColor(0xed, 0xc9, 0x3f),
}


def _load_data(args: argparse.Namespace) -> dict:
    if args.data and args.data_file:
        raise SystemExit("Use either --data or --data-file (not both)")

    if args.data_file:
        p = Path(args.data_file)
        return json.loads(p.read_text(encoding="utf-8"))

    if args.data:
        return json.loads(args.data)

    return {}

def add_gradient_background(slide, color1, color2):
    """グラデーション背景を追加"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # グラデーション効果はXMLで設定可能（簡略化）
    return background

def add_premium_title_slide(prs):
    """プレミアムタイトルスライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # 装飾円
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['secondary']
    circle1.fill.fore_color.brightness = 0.3
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent']
    circle2.fill.fore_color.brightness = -0.2
    circle2.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Transformation"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Calibri Light"

    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategic Roadmap 2026"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.font.name = "Calibri Light"

    # 区切り線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()

    # メタ情報
    meta_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.8))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Confidential | Board of Directors | January 2026"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['white']
    p.font.color.brightness = 0.3

    return slide

def add_infographic_slide(prs):
    """インフォグラフィックスライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light']
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Performance Indicators"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # 4つのKPIカード
    kpis = [
        ('¥12.5B', 'Revenue', '+23%', COLORS['success']),
        ('18.5%', 'Margin', '+2.3pp', COLORS['secondary']),
        ('2,450', 'Employees', '+12%', COLORS['accent']),
        ('94%', 'NPS Score', '+5pp', COLORS['accent2']),
    ]

    for i, (value, label, change, color) in enumerate(kpis):
        x = Inches(0.5 + i * 3.2)
        y = Inches(1.5)

        # カード背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.shadow.inherit = False

        # アイコン円
        icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y + Inches(0.2), Inches(0.7), Inches(0.7))
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = color
        icon_circle.line.fill.background()

        # 値
        value_box = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(2.9), Inches(0.6))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

        # ラベル
        label_box = slide.shapes.add_textbox(x, y + Inches(1.55), Inches(2.9), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # 変化率
        change_box = slide.shapes.add_textbox(x, y + Inches(1.9), Inches(2.9), Inches(0.3))
        tf = change_box.text_frame
        p = tf.paragraphs[0]
        p.text = change
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # 下部チャートエリア
    chart_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(12.333), Inches(3.2))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = COLORS['white']
    chart_bg.line.fill.background()

    # チャートタイトル
    chart_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(5), Inches(0.5))
    tf = chart_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Trend"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 簡易チャート（バー）
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    values = [8.5, 9.2, 10.1, 11.0, 11.8, 12.5]
    max_val = max(values)

    bar_width = Inches(0.6)
    spacing = Inches(1.8)
    start_x = Inches(1.5)
    base_y = Inches(6.8)

    for i, (month, val) in enumerate(zip(months, values)):
        bar_height = Inches(2.2 * val / max_val)
        x = start_x + i * spacing
        y = base_y - bar_height

        # バー
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS['primary'] if i < 3 else COLORS['accent']
        bar.line.fill.background()

        # 月ラベル
        month_box = slide.shapes.add_textbox(x - Inches(0.1), base_y, Inches(0.8), Inches(0.3))
        tf = month_box.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # 値ラベル
        val_box = slide.shapes.add_textbox(x - Inches(0.1), y - Inches(0.25), Inches(0.8), Inches(0.3))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'¥{val}B'
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_process_diagram_slide(prs):
    """プロセス図スライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Digital Transformation Roadmap"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # プロセスステップ
    steps = [
        ('Strategy', 'Define vision\n& goals', COLORS['primary']),
        ('Assessment', 'Current state\nanalysis', COLORS['secondary']),
        ('Design', 'Solution\narchitecture', COLORS['accent']),
        ('Implementation', 'Deploy &\nintegrate', COLORS['success']),
        ('Optimization', 'Monitor &\nimprove', COLORS['accent2']),
    ]

    start_x = Inches(0.8)
    y = Inches(3)
    box_width = Inches(2.2)
    box_height = Inches(2)
    spacing = Inches(2.6)

    for i, (title, desc, color) in enumerate(steps):
        x = start_x + i * spacing

        # ステップボックス
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # ステップ番号
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y - Inches(0.4), Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLORS['white']
        num_circle.line.color.rgb = color
        num_circle.line.width = Pt(2)

        num_box = slide.shapes.add_textbox(x + Inches(0.8), y - Inches(0.35), Inches(0.6), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # タイトル
        title_tb = slide.shapes.add_textbox(x, y + Inches(0.3), box_width, Inches(0.5))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 説明
        desc_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), Inches(2.0), Inches(0.8))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 接続矢印
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x + box_width + Inches(0.05), y + Inches(0.6), Inches(0.4), Inches(0.8))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gray']
            arrow.fill.fore_color.brightness = 0.5
            arrow.line.fill.background()

    return slide

def add_swot_slide(prs):
    """SWOT分析スライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light']
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SWOT Analysis"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # SWOTグリッド
    swot_data = [
        ('Strengths', 'S', COLORS['success'], [
            'Strong brand recognition',
            'Innovative product pipeline',
            'High customer retention'
        ]),
        ('Weaknesses', 'W', COLORS['accent2'], [
            'Limited market presence',
            'Legacy infrastructure',
            'Resource constraints'
        ]),
        ('Opportunities', 'O', COLORS['secondary'], [
            'Emerging markets growth',
            'Digital transformation',
            'Strategic partnerships'
        ]),
        ('Threats', 'T', COLORS['gray'], [
            'Intense competition',
            'Regulatory changes',
            'Economic uncertainty'
        ]),
    ]

    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(6.7), Inches(1.3)),
        (Inches(0.5), Inches(4.4)),
        (Inches(6.7), Inches(4.4)),
    ]

    for (title, letter, color, items), (x, y) in zip(swot_data, positions):
        # ボックス
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.133), Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # ヘッダー背景
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(6.133), Inches(0.6))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()

        # タイトル
        title_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5), Inches(0.4))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']

        # リストアイテム
        for j, item in enumerate(items):
            item_y = y + Inches(0.8 + j * 0.55)

            # バレット
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), item_y + Inches(0.1), Inches(0.12), Inches(0.12))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = color
            bullet.line.fill.background()

            # テキスト
            item_tb = slide.shapes.add_textbox(x + Inches(0.55), item_y, Inches(5.3), Inches(0.5))
            tf = item_tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark']

    return slide

def add_timeline_slide(prs):
    """タイムラインスライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Implementation Timeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # タイムラインライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4), Inches(11.733), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.fill.fore_color.brightness = 0.5
    line.line.fill.background()

    # マイルストーン
    milestones = [
        ('Q1 2026', 'Planning', 'Requirements\nDocumentation'),
        ('Q2 2026', 'Design', 'Architecture\nApproval'),
        ('Q3 2026', 'Development', 'Core Features\nComplete'),
        ('Q4 2026', 'Launch', 'Production\nDeployment'),
    ]

    for i, (quarter, phase, detail) in enumerate(milestones):
        x = Inches(1.5 + i * 3)

        # マイルストーン円
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(3.7), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent'] if i == 0 else COLORS['primary']
        circle.line.color.rgb = COLORS['white']
        circle.line.width = Pt(3)

        # 四半期
        quarter_tb = slide.shapes.add_textbox(x - Inches(0.5), Inches(2.5), Inches(1.6), Inches(0.4))
        tf = quarter_tb.text_frame
        p = tf.paragraphs[0]
        p.text = quarter
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent'] if i == 0 else COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

        # フェーズ
        phase_tb = slide.shapes.add_textbox(x - Inches(0.5), Inches(2.9), Inches(1.6), Inches(0.4))
        tf = phase_tb.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # 詳細（下）
        detail_tb = slide.shapes.add_textbox(x - Inches(0.6), Inches(4.5), Inches(1.8), Inches(0.8))
        tf = detail_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

        # 期間ラベル
        if i < len(milestones) - 1:
            duration_tb = slide.shapes.add_textbox(x + Inches(0.8), Inches(3.85), Inches(2), Inches(0.3))
            tf = duration_tb.text_frame
            p = tf.paragraphs[0]
            p.text = '3 months'
            p.font.size = Pt(9)
            p.font.color.rgb = COLORS['gray']
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_team_slide(prs):
    """チーム紹介スライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # チームメンバー
    team = [
        ('CEO', 'Yamada Taro', 'Vision & Strategy'),
        ('CTO', 'Sato Kenji', 'Technology & Innovation'),
        ('CFO', 'Tanaka Yuki', 'Finance & Operations'),
        ('COO', 'Suzuki Akiko', 'Business Development'),
    ]

    for i, (role, name, focus) in enumerate(team):
        x = Inches(0.5 + i * 3.2)
        y = Inches(2)

        # プロフィールカード
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.fill.background()

        # アバター円（プレースホルダー）
        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.4), Inches(1.2), Inches(1.2))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = COLORS['light']
        avatar.line.color.rgb = COLORS['accent']
        avatar.line.width = Pt(3)

        # イニシャル
        initial = name.split()[0][0] + name.split()[1][0]
        init_tb = slide.shapes.add_textbox(x + Inches(0.85), y + Inches(0.8), Inches(1.2), Inches(0.5))
        tf = init_tb.text_frame
        p = tf.paragraphs[0]
        p.text = initial
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # 名前
        name_tb = slide.shapes.add_textbox(x, y + Inches(1.8), Inches(2.9), Inches(0.5))
        tf = name_tb.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

        # 役職
        role_tb = slide.shapes.add_textbox(x, y + Inches(2.3), Inches(2.9), Inches(0.4))
        tf = role_tb.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['accent']
        p.alignment = PP_ALIGN.CENTER

        # フォーカス
        focus_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(2.9), Inches(2.5), Inches(1.2))
        tf = focus_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = focus
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_advanced_presentation(output_path: str, _data: dict | None = None):
    """高度なビジネスプレゼンテーションを作成"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # スライド追加
    add_premium_title_slide(prs)
    add_infographic_slide(prs)
    add_process_diagram_slide(prs)
    add_swot_slide(prs)
    add_timeline_slide(prs)
    add_team_slide(prs)

    # 保存
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate advanced business PPTX (python-pptx)")
    parser.add_argument("--output", default=str(Path("output/pptx/06_advanced_business.pptx")))
    parser.add_argument("--data", help="Inline JSON string (optional)")
    parser.add_argument("--data-file", help="JSON file path (optional)")
    args = parser.parse_args()

    data = _load_data(args)
    create_advanced_presentation(args.output, data)
